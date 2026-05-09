"""
Генератор презентации DFORM Studio — дипломная работа
Северо-Кавказская государственная академия, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── Цветовая палитра ───────────────────────────────────────────────────────
BG        = RGBColor(0x0C, 0x0C, 0x0E)
BG_CARD   = RGBColor(0x18, 0x18, 0x1B)
BG_CARD2  = RGBColor(0x22, 0x22, 0x26)
BG_FIELD  = RGBColor(0x14, 0x14, 0x17)   # фон строки поля в таблице БД
ACCENT    = RGBColor(0xE8, 0xF4, 0x39)
TEXT      = RGBColor(0xF4, 0xEF, 0xE6)
TEXT_MUTED= RGBColor(0x88, 0x85, 0x80)
SEPARATOR = RGBColor(0x2A, 0x2A, 0x2E)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

IMG_DIR = r'd:\Diploms\2025-2026\diiinav\project\screenshots_from_docx'

# ════════════════════════════════════════════════════════════════════════════
# ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ
# ════════════════════════════════════════════════════════════════════════════

def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h,
       size=16, bold=False, italic=False,
       color=TEXT, align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.word_wrap = wrap
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.333, 0.06, ACCENT)
    tb(slide, title, 0.5, 0.14, 11.0, 0.72, size=28, bold=True, color=TEXT)
    if subtitle:
        tb(slide, subtitle, 0.5, 0.82, 11.0, 0.38, size=13, color=TEXT_MUTED)
    rect(slide, 0, 7.38, 13.333, 0.06, ACCENT)
    tb(slide, "DFORM Studio", 11.3, 7.2, 1.9, 0.22,
       size=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def insert_image(slide, img_path, x, y, w, h):
    """Вставить изображение с сохранением пропорций."""
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
        return True
    return False


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛЬНЫЙ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()

rect(slide, 0, 0, 0.45, 7.5, ACCENT)
rect(slide, 0.45, 0, 13.0, 0.06, ACCENT)
rect(slide, 0.45, 7.44, 13.0, 0.06, ACCENT)

rect(slide, 1.0, 0.55, 4.0, 0.42, ACCENT)
tb(slide, "ДИПЛОМНАЯ РАБОТА", 1.05, 0.6, 3.9, 0.33, size=12, bold=True, color=BG)

# Название
tb(slide, "Разработка корпоративного сайта", 1.0, 1.2, 10.5, 0.75,
   size=34, bold=True, color=TEXT)
tb(slide, "с возможностью онлайн-заказов", 1.0, 1.92, 10.5, 0.7,
   size=34, bold=True, color=ACCENT)
tb(slide, "и аналитикой", 1.0, 2.62, 10.5, 0.65,
   size=34, bold=True, color=TEXT)

rect(slide, 1.0, 3.45, 11.3, 0.04, SEPARATOR)

rows = [
    ("Студент:",       ""),
    ("Направление:",   "09.03.03 — Прикладная информатика"),
    ("Руководитель:",  "(ФИО руководителя)"),
    ("Кафедра:",       "Информационных технологий"),
]
for i, (label, val) in enumerate(rows):
    cy = 3.62 + i * 0.52
    tb(slide, label, 1.0,  cy, 2.4, 0.42, size=12, color=TEXT_MUTED)
    tb(slide, val,   3.45, cy, 8.0, 0.42, size=14, bold=True, color=TEXT)

rect(slide, 1.0, 5.95, 11.3, 0.04, SEPARATOR)
tb(slide, "Северо-Кавказская государственная академия  ·  Черкесск, 2026",
   1.0, 6.12, 11.3, 0.38, size=13, color=TEXT_MUTED)

# Декоративный блок
rect(slide, 10.6, 1.1, 2.15, 2.15, BG_CARD)
rect(slide, 10.6, 1.1, 2.15, 0.06, ACCENT)
tb(slide, "DF", 10.6, 1.25, 2.15, 1.4, size=60, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(slide, "DFORM Studio", 10.6, 2.75, 2.15, 0.35, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — АКТУАЛЬНОСТЬ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Актуальность темы")

cards = [
    ("70%+",              "потребителей изучают сайт\nперед обращением к компании"),
    ("Онлайн-заказы",     "снимают ограничение рабочих часов\nи структурируют первичные брифы"),
    ("Аналитика данных",  "необходима для обоснованного\nпланирования загрузки и услуг"),
    ("Готовые платформы", "не покрывают специфику\nтворческих компаний"),
    ("CMS-панель",        "обеспечивает независимость\nот разработчика при обновлениях"),
]

for i, (title, desc) in enumerate(cards):
    col = i % 3
    row = i // 3
    cx = 0.4 + col * 4.35
    cy = 1.22 + row * 2.55
    rect(slide, cx, cy, 4.1, 2.3, BG_CARD)
    rect(slide, cx, cy, 4.1, 0.06, ACCENT)
    tb(slide, title, cx+0.2, cy+0.18, 3.7, 0.52, size=18, bold=True, color=ACCENT)
    tb(slide, desc,  cx+0.2, cy+0.72, 3.7, 1.3,  size=13, color=TEXT)

rect(slide, 0.4, 6.5, 12.5, 0.58, BG_CARD2)
tb(slide, "⟶  Индивидуальная разработка — единственное решение, обеспечивающее "
          "полный функционал без компромиссов с дизайном, гибкостью и независимостью",
   0.65, 6.55, 12.0, 0.48, size=12, italic=True, color=TEXT_MUTED)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — ОБЪЕКТ / ПРЕДМЕТ / ЦЕЛЬ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Объект, предмет и цель работы")

blocks = [
    ("Объект исследования",
     "Процесс автоматизации взаимодействия компании сферы услуг с клиентами посредством корпоративного веб-сайта",
     1.22),
    ("Предмет исследования",
     "Методы и средства разработки корпоративного сайта с интегрированными модулями онлайн-заказов и аналитики на основе современного стека веб-технологий",
     2.97),
    ("Цель работы",
     "Проектирование и реализация корпоративного сайта, обеспечивающего представление информации об услугах, приём и управление онлайн-заявками, а также сбор и визуализацию аналитических данных о деятельности студии",
     4.75),
]

for label, content, cy in blocks:
    rect(slide, 0.4, cy, 12.5, 1.45, BG_CARD)
    rect(slide, 0.4, cy, 0.07, 1.45, ACCENT)
    tb(slide, label,   0.65, cy+0.1,  3.8, 0.38, size=12, bold=True, color=ACCENT)
    tb(slide, content, 0.65, cy+0.52, 11.9, 0.8, size=13.5, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — ЗАДАЧИ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Задачи дипломной работы")

tasks = [
    "Анализ предметной области и потребностей компаний сферы услуг",
    "Обзор существующих решений и обоснование индивидуальной разработки",
    "Формулирование функциональных и нефункциональных требований",
    "Выбор и обоснование технологического стека",
    "Проектирование структуры БД и архитектуры приложения",
    "Реализация серверной части с REST API и аутентификацией",
    "Разработка клиентского интерфейса с адаптивной вёрсткой",
    "Реализация системы заказов, CMS-панели и модуля аналитики",
    "Функциональное тестирование и развёртывание на Render.com",
]

col1 = tasks[:5]
col2 = tasks[5:]

for i, task in enumerate(col1):
    cy = 1.2 + i * 1.06
    rect(slide, 0.4, cy, 0.44, 0.44, ACCENT)
    tb(slide, str(i+1), 0.4, cy+0.04, 0.44, 0.38, size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(slide, 0.91, cy, 5.85, 0.44, BG_CARD)
    tb(slide, task, 1.08, cy+0.08, 5.6, 0.32, size=13, color=TEXT)

for i, task in enumerate(col2):
    cy = 1.2 + i * 1.06
    n = i + 6
    rect(slide, 7.1, cy, 0.44, 0.44, ACCENT)
    tb(slide, str(n), 7.1, cy+0.04, 0.44, 0.38, size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(slide, 7.61, cy, 5.35, 0.44, BG_CARD)
    tb(slide, task, 7.78, cy+0.08, 5.1, 0.32, size=13, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — СТЕК ТЕХНОЛОГИЙ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Технологический стек")

tech_groups = [
    ("Фронтенд",      ["React 19", "React Router 7", "Context API + Хуки", "CSS Custom Properties", "Fetch API"]),
    ("Бэкенд",        ["Node.js", "Встроенный http-сервер", "REST API", "bcrypt (хэширование)", "Сессии + Cookies"]),
    ("База данных",   ["MySQL", "16 таблиц", "5 внешних ключей", "Нормализация 3НФ", "schema.sql"]),
    ("Инфраструктура",["Git + GitHub", "Render.com", "CI/CD авто-деплой", "HTTPS / TLS", "ENV переменные"]),
]

for i, (group, items) in enumerate(tech_groups):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * 6.55
    cy = 1.18 + row * 2.9
    rect(slide, cx, cy, 6.2, 2.68, BG_CARD)
    rect(slide, cx, cy, 6.2, 0.06, ACCENT)
    tb(slide, group, cx+0.22, cy+0.14, 5.8, 0.48, size=16, bold=True, color=ACCENT)
    for j, item in enumerate(items):
        tb(slide, f"  {item}", cx+0.22, cy+0.65 + j*0.38, 5.8, 0.36, size=14, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — АРХИТЕКТУРА
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Архитектура системы")

components = [
    (0.5,  2.1, 3.1, 0.95, "БРАУЗЕР",        "React SPA",        ACCENT,  BG),
    (4.4,  2.1, 3.1, 0.95, "NODE.JS СЕРВЕР", "REST API",         BG_CARD2, ACCENT),
    (8.3,  2.1, 3.1, 0.95, "JSON / MySQL",   "16 таблиц",        BG_CARD,  TEXT),
    (0.5,  4.0, 3.1, 0.95, "ПОСЕТИТЕЛЬ",     "Публичный сайт",   BG_CARD,  TEXT_MUTED),
    (4.4,  4.0, 3.1, 0.95, "АДМИНИСТРАТОР",  "CMS + Аналитика",  BG_CARD,  TEXT_MUTED),
    (8.3,  4.0, 3.1, 0.95, "RENDER.COM",     "Облачный хостинг", BG_CARD,  TEXT_MUTED),
]

for (cx, cy, w, h, line1, line2, bg_c, txt_c) in components:
    rect(slide, cx, cy, w, h, bg_c)
    tb(slide, line1, cx, cy+0.07, w, 0.42, size=13, bold=True, color=txt_c, align=PP_ALIGN.CENTER)
    sub_c = BG if bg_c == ACCENT else (ACCENT if bg_c == BG_CARD2 else TEXT_MUTED)
    tb(slide, line2, cx, cy+0.52, w, 0.35, size=11, color=sub_c, align=PP_ALIGN.CENTER)

def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = ACCENT
    conn.line.width = Pt(1.5)

arrow(slide, 3.6, 2.58, 4.4, 2.58)
arrow(slide, 7.5, 2.58, 8.3, 2.58)
arrow(slide, 2.05, 3.05, 2.05, 4.0)
arrow(slide, 6.0,  3.05, 6.0,  4.0)
arrow(slide, 9.85, 3.05, 9.85, 4.0)

tb(slide, "HTTP/JSON", 3.65, 2.28, 1.4, 0.3, size=9, color=TEXT_MUTED)
tb(slide, "fs.read/write", 7.55, 2.28, 1.5, 0.3, size=9, color=TEXT_MUTED)

rect(slide, 0.5, 5.5, 12.0, 0.82, BG_CARD)
tb(slide, "Клиент-серверная SPA-архитектура: фронтенд (React) ↔ бэкенд (Node.js HTTP) ↔ хранилище (MySQL / JSON).\n"
          "В production: Node.js раздаёт скомпилированный React-билд как статику — единый сервис на Render.",
   0.72, 5.56, 11.56, 0.7, size=12, color=TEXT_MUTED)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — СХЕМА БД (поля вертикально, стиль DrawSQL)
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Структура базы данных (MySQL)", "16 таблиц · 5 внешних ключей · Нормализация до 3НФ")

# Таблицы: (название, [(поле, тип, флаги)])
tables_db = [
    ("orders", [
        ("id",          "VARCHAR(64)", "PK"),
        ("name",        "VARCHAR(255)", ""),
        ("email",       "VARCHAR(255)", ""),
        ("phone",       "VARCHAR(50)",  "NULL"),
        ("service_id",  "VARCHAR(64)",  "FK"),
        ("description", "TEXT",         ""),
        ("budget",      "VARCHAR(100)", "NULL"),
        ("status",      "ENUM",         "new|in_prog…"),
        ("created_at",  "DATETIME",     ""),
    ]),
    ("services", [
        ("id",          "VARCHAR(64)",  "PK"),
        ("icon",        "VARCHAR(10)",  ""),
        ("title",       "VARCHAR(255)", ""),
        ("description", "TEXT",         ""),
        ("price",       "VARCHAR(100)", ""),
        ("duration",    "VARCHAR(100)", ""),
    ]),
    ("service_features", [
        ("id",          "INT",         "PK AI"),
        ("service_id",  "VARCHAR(64)", "FK"),
        ("feature",     "VARCHAR(255)",""),
        ("sort_order",  "INT",         ""),
    ]),
    ("portfolio", [
        ("id",          "VARCHAR(64)",  "PK"),
        ("title",       "VARCHAR(255)", ""),
        ("category",    "VARCHAR(100)", ""),
        ("client",      "VARCHAR(255)", ""),
        ("year",        "VARCHAR(10)",  ""),
        ("description", "TEXT",         ""),
        ("color",       "VARCHAR(20)",  ""),
        ("stats_*",     "VARCHAR(100)", "×3"),
    ]),
    ("portfolio_tags", [
        ("id",           "INT",         "PK AI"),
        ("portfolio_id", "VARCHAR(64)", "FK"),
        ("tag",          "VARCHAR(100)",""),
        ("sort_order",   "INT",         ""),
    ]),
    ("team", [
        ("id",       "INT",          "PK AI"),
        ("name",     "VARCHAR(255)", ""),
        ("role",     "VARCHAR(255)", ""),
        ("bio",      "TEXT",         ""),
        ("initials", "VARCHAR(5)",   ""),
        ("color",    "VARCHAR(20)",  ""),
    ]),
    ("analytics", [
        ("id",            "INT", "PK AI"),
        ("total_orders",  "INT", ""),
        ("total_revenue", "INT", ""),
    ]),
    ("analytics_orders_by_day", [
        ("id",           "INT",  "PK AI"),
        ("analytics_id", "INT",  "FK"),
        ("date",         "DATE", "UNIQUE"),
        ("count",        "INT",  ""),
    ]),
]

# 4 колонки × 2 строки
col_w   = 3.1
row_h_base = 0.32   # высота строки поля
header_h   = 0.38   # высота заголовка таблицы
padding_v  = 0.1
cols_count = 4

for i, (tname, fields) in enumerate(tables_db):
    col = i % cols_count
    row = i // cols_count
    cx = 0.22 + col * (col_w + 0.12)
    body_h = header_h + padding_v + len(fields) * row_h_base + padding_v
    cy = 1.15 + row * 3.4

    # Фон таблицы
    rect(slide, cx, cy, col_w, body_h, BG_CARD)
    # Заголовок таблицы
    rect(slide, cx, cy, col_w, header_h, BG_CARD2)
    rect(slide, cx, cy, 0.06, body_h, ACCENT)   # левая акцентная полоска
    tb(slide, tname, cx+0.12, cy+0.06, col_w-0.18, 0.28, size=11, bold=True, color=ACCENT)

    # Поля
    for j, (fname, ftype, flag) in enumerate(fields):
        fy = cy + header_h + padding_v + j * row_h_base
        # Чередующийся фон строк
        if j % 2 == 0:
            rect(slide, cx+0.06, fy, col_w-0.06, row_h_base, BG_FIELD)
        # Имя поля
        field_color = ACCENT if flag in ("PK", "PK AI") else (TEXT_MUTED if flag == "FK" else TEXT)
        tb(slide, fname, cx+0.1, fy+0.02, 1.35, row_h_base-0.04, size=8.5, color=field_color)
        # Тип
        tb(slide, ftype, cx+1.48, fy+0.02, 1.1, row_h_base-0.04, size=7.5, color=TEXT_MUTED)
        # Флаг
        if flag:
            flag_color = ACCENT if "PK" in flag else (RGBColor(0x60,0x9A,0xF5) if flag=="FK" else TEXT_MUTED)
            tb(slide, flag, cx+2.58, fy+0.02, 0.5, row_h_base-0.04, size=7, bold=True, color=flag_color)

# Нижняя подпись
tb(slide, "PK = первичный ключ  ·  FK = внешний ключ  ·  AI = AUTO_INCREMENT  ·  "
          "Не показано: site_hero, site_about, site_cta, site_footer, clients, stats, company_values, contacts, social_links",
   0.22, 7.08, 12.8, 0.28, size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙДЫ 8–16 — ЭКРАННЫЕ ФОРМЫ
# ════════════════════════════════════════════════════════════════════════════

# Маппинг: (заголовок, подзаголовок, файл_изображения, подпись_рисунка, ключевые_особенности)
screens = [
    (
        "1. Главная страница",
        "Hero-блок, счётчики, портфолио preview, услуги preview, CTA",
        "image1.png",
        "Рисунок 1 — Главная страница",
        ["Полноэкранный hero с заголовком", "Счётчики (150+ проектов, 5 лет)", "Превью портфолио (3 кейса)", "Превью услуг (карточки)", "CTA-блок с кнопкой заказа"],
    ),
    (
        "2. Страница портфолио",
        "Сетка проектов с фильтрацией по категориям",
        "image2.png",
        "Рисунок 2 — Страница портфолио с фильтрацией",
        ["Фильтрация по категориям", "Карточки с акцентным цветом", "Детальная страница кейса", "Теги и статистика проекта", "Клиентская фильтрация (без запросов)"],
    ),
    (
        "3. Страница услуг",
        "Карточки услуг с ценами и перечнем работ",
        "image3.png",
        "Рисунок 3 — Страница услуг",
        ["Иконка + название услуги", "Описание и ценовой ориентир", "Срок выполнения", "Перечень включённых работ", "Кнопка оформления заказа"],
    ),
    (
        "4. Страница «О студии»",
        "История, команда, ценности, клиенты, статистика",
        "image4.png",
        "Рисунок 4 — Страница «О студии»",
        ["История студии (текст из CMS)", "Карточки команды с аватарами", "Блок ценностей компании", "Список клиентов студии", "Числовые показатели"],
    ),
    (
        "5. Форма заказа",
        "Многошаговая форма: услуга → описание → бюджет → контакты",
        "image6.png",
        "Рисунок 6 — Форма заказа, шаг 1 — выбор услуги",
        ["Шаг 1: выбор услуги (карточки)", "Шаг 2: описание задачи", "Шаг 3: бюджет (опционально)", "Шаг 4: контактные данные", "Индикатор прогресса + валидация"],
    ),
    (
        "6. Административная панель — Заказы",
        "Список заявок, детальное модальное окно, смена статуса",
        "image9.png",
        "Рисунок 9 — Список заявок в административной панели",
        ["Таблица заявок (новые сверху)", "Цветовое кодирование статусов", "Модальное окно с деталями", "Активные ссылки: mailto + tel", "Смена статуса без перезагрузки"],
    ),
    (
        "7. Административная панель — CMS",
        "Управление контентом всех разделов сайта",
        "image12.png",
        "Рисунок 12 — Форма редактирования проекта в портфолио",
        ["Редактирование портфолио (+ теги)", "Управление услугами (+ список работ)", "Редактирование команды (аватар)", "9 подразделов настроек сайта", "Изменения → мгновенно на сайте"],
    ),
    (
        "8. Модуль аналитики",
        "Сводные показатели, график, распределение по статусам и услугам",
        "image11.png",
        "Рисунок 11 — Аналитический дашборд",
        ["Карточки: всего / новых / в работе", "Линейный SVG-график (30 дней)", "Распределение по статусам", "Распределение по услугам", "Кнопка обновления данных"],
    ),
    (
        "9. Адаптивная вёрстка",
        "Корректное отображение от 320px до 1920px на всех устройствах",
        "image13.png",
        "Рисунок 13 — Сравнение десктопного и мобильного вида",
        ["CSS media queries (768px, 480px)", "Мобильное меню навигации", "Одноколоночные сетки на mobile", "clamp() для масштабирования шрифтов", "Протестировано: iPhone SE, iPad, FHD"],
    ),
]

for idx, (title, subtitle, img_file, fig_caption, points) in enumerate(screens):
    slide = add_slide()
    slide_header(slide, f"Экранные формы — {title}", subtitle)

    img_path = os.path.join(IMG_DIR, img_file)
    img_x, img_y, img_w, img_h = 0.35, 1.15, 7.6, 5.85

    # Левый блок — изображение или заглушка
    rect(slide, img_x, img_y, img_w, img_h, BG_CARD)
    rect(slide, img_x, img_y, img_w, 0.05, ACCENT)

    if os.path.exists(img_path):
        # Вставляем изображение поверх фона с небольшим отступом
        slide.shapes.add_picture(img_path,
                                  Inches(img_x + 0.08), Inches(img_y + 0.08),
                                  Inches(img_w - 0.16),  Inches(img_h - 0.28))
    else:
        tb(slide, "[ Скриншот ]", img_x, img_y + img_h/2 - 0.3, img_w, 0.6,
           size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Подпись рисунка
    tb(slide, fig_caption,
       img_x, img_y + img_h - 0.18, img_w, 0.22,
       size=8.5, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Правый блок — ключевые особенности
    rx = 8.22
    rect(slide, rx, 1.15, 4.78, 5.85, BG_CARD2)
    rect(slide, rx, 1.15, 4.78, 0.05, ACCENT)
    tb(slide, "Ключевые особенности", rx+0.2, 1.27, 4.35, 0.4, size=13, bold=True, color=ACCENT)

    for j, point in enumerate(points):
        py = 1.82 + j * 0.96
        rect(slide, rx, py, 0.07, 0.6, ACCENT)
        tb(slide, point, rx+0.22, py+0.08, 4.4, 0.5, size=13, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 17 — РЕЗУЛЬТАТЫ ТЕСТИРОВАНИЯ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Результаты тестирования")

test_results = [
    ("Функциональное\nтестирование",   "Все ключевые сценарии пройдены:\nзаказ, аутентификация, CMS, аналитика", "✓  Пройдено"),
    ("Кросс-браузерное\nтестирование", "Chrome, Firefox, Edge, Safari —\nвсе браузеры работают корректно",       "✓  Пройдено"),
    ("Адаптивность\n320 – 1920px",     "iPhone SE, iPhone 12, iPad, Full HD —\nбез замечаний ко всем форматам",   "✓  Пройдено"),
    ("Безопасность\nадмин-панели",     "401 без сессии, bcrypt-пароли,\nHttpOnly cookies, CORS настроен",         "✓  Пройдено"),
    ("Производительность\nAPI",        "Время ответа API < 200 мс,\nзагрузка страниц < 2 с",                     "✓  Норма"),
    ("Выявленные дефекты",             "3 дефекта обнаружено и устранено:\nSVG/Safari, CORS, service_features",   "→  Исправлено"),
]

for i, (label, desc, status) in enumerate(test_results):
    col = i % 3
    row = i // 3
    cx = 0.4 + col * 4.35
    cy = 1.22 + row * 2.65
    rect(slide, cx, cy, 4.1, 2.45, BG_CARD)
    rect(slide, cx, cy, 4.1, 0.06, ACCENT)
    tb(slide, label,  cx+0.2, cy+0.18, 3.7, 0.62, size=14, bold=True, color=ACCENT)
    tb(slide, desc,   cx+0.2, cy+0.82, 3.7, 0.98, size=12, color=TEXT)
    tb(slide, status, cx+0.2, cy+1.98, 3.7, 0.36, size=12, bold=True, color=ACCENT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 18 — ВЫВОДЫ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()
slide_header(slide, "Выводы")

conclusions = [
    ("Цель достигнута",                "Корпоративный сайт с онлайн-заказами, CMS и аналитикой спроектирован и реализован"),
    ("Система готова к эксплуатации",  "Приложение развёрнуто на Render.com и доступно по публичному URL"),
    ("Независимость от разработчика",  "Административная панель управляет 9 разделами сайта без знания кода"),
    ("Информационная основа решений",  "Аналитический дашборд с графиком и распределениями обеспечивает данные для планирования"),
    ("Архитектура масштабируема",      "Решение адаптируемо для других компаний сферы услуг с минимальными доработками"),
    ("Требования выполнены",           "Все функциональные и нефункциональные требования подтверждены тестированием"),
]

for i, (title, desc) in enumerate(conclusions):
    cy = 1.22 + i * 1.0
    rect(slide, 0.4, cy, 0.46, 0.74, ACCENT)
    tb(slide, str(i+1), 0.4, cy+0.14, 0.46, 0.5, size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    rect(slide, 0.93, cy, 12.0, 0.74, BG_CARD)
    tb(slide, title, 1.12, cy+0.06, 4.8, 0.36, size=14, bold=True, color=ACCENT)
    tb(slide, desc,  1.12, cy+0.4,  11.65, 0.3, size=12.5, color=TEXT)


# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 19 — ФИНАЛЬНЫЙ
# ════════════════════════════════════════════════════════════════════════════
slide = add_slide()

rect(slide, 0, 0, 0.45, 7.5, ACCENT)
rect(slide, 0.45, 0, 13.0, 0.07, ACCENT)
rect(slide, 0.45, 7.43, 13.0, 0.07, ACCENT)

tb(slide, "Доклад окончен.", 1.0, 1.7, 11.5, 1.35, size=58, bold=True, color=TEXT)
tb(slide, "Спасибо за внимание!", 1.0, 3.1, 11.5, 0.95, size=36, color=ACCENT)

rect(slide, 1.0, 4.35, 10.5, 0.04, SEPARATOR)

tb(slide, "Северо-Кавказская государственная академия  ·  Черкесск, 2026",
   1.0, 4.55, 10.5, 0.42, size=14, color=TEXT_MUTED)

rect(slide, 10.6, 4.3, 2.3, 2.3, BG_CARD)
rect(slide, 10.6, 4.3, 2.3, 0.06, ACCENT)
tb(slide, "DF", 10.6, 4.45, 2.3, 1.4, size=64, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(slide, "DFORM Studio", 10.6, 5.9, 2.3, 0.38, size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ─── Сохранение ─────────────────────────────────────────────────────────────
out_path = r"d:\Diploms\2025-2026\diiinav\project\Презентация_DFORM.pptx"
prs.save(out_path)
print(f"Сохранено: {out_path}")
print(f"Слайдов: {len(prs.slides)}")
